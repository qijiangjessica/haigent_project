"""
Haigent Project — PPTX Generator
Run: pip install python-pptx && python generate_pptx.py
Output: Haigent_Project.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand colours ─────────────────────────────────────────────────────────────
CREAM      = RGBColor(0xF8, 0xF4, 0xED)
CHARCOAL   = RGBColor(0x23, 0x23, 0x23)
GOLD       = RGBColor(0xF3, 0xCF, 0x63)
TEAL       = RGBColor(0x19, 0xA9, 0xB6)
RED        = RGBColor(0xE3, 0x5B, 0x6D)
GREEN      = RGBColor(0x9A, 0xBF, 0x45)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xED, 0xE8, 0xDF)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # completely blank layout


# ── Helpers ───────────────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width or 1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=CHARCOAL,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_multiline(slide, lines, left, top, width, height,
                  font_size=14, color=CHARCOAL, line_spacing=1.2):
    """lines = list of (text, bold, color_override)"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            text, bold, col = item, False, color
        else:
            text = item[0]
            bold = item[1] if len(item) > 1 else False
            col  = item[2] if len(item) > 2 else color

        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = col
    return txBox


def bg(slide, color=CREAM):
    add_rect(slide, 0, 0, 13.33, 7.5, color)


def accent_bar(slide, color=GOLD, height=0.08):
    add_rect(slide, 0, 0, 13.33, height, color)


def slide_header(slide, title, subtitle=None, accent=GOLD):
    bg(slide)
    accent_bar(slide, accent)
    add_text(slide, title, 0.5, 0.2, 12, 0.7,
             font_size=32, bold=True, color=CHARCOAL)
    if subtitle:
        add_text(slide, subtitle, 0.5, 0.9, 12, 0.5,
                 font_size=16, color=RGBColor(0x6B, 0x65, 0x60))


def pill(slide, text, left, top, width=1.8, height=0.38,
         bg_color=GOLD, text_color=CHARCOAL, font_size=12):
    add_rect(slide, left, top, width, height, bg_color)
    add_text(slide, text, left + 0.05, top + 0.04, width - 0.1, height - 0.08,
             font_size=font_size, bold=True, color=text_color, align=PP_ALIGN.CENTER)


def table_row(slide, cols, left, top, col_widths, row_height=0.42,
              bg_color=WHITE, text_color=CHARCOAL, font_size=12, bold=False):
    x = left
    for i, (text, w) in enumerate(zip(cols, col_widths)):
        add_rect(slide, x, top, w, row_height, bg_color,
                 line_color=LIGHT_GREY, line_width=0.5)
        add_text(slide, text, x + 0.08, top + 0.06, w - 0.16, row_height - 0.1,
                 font_size=font_size, bold=bold, color=text_color)
        x += w


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, CHARCOAL)

# Gold top stripe
add_rect(sl, 0, 0, 13.33, 0.12, GOLD)
# Gold bottom stripe
add_rect(sl, 0, 7.38, 13.33, 0.12, GOLD)

# Main title
add_text(sl, "Haigent", 1, 1.6, 11, 1.4,
         font_size=72, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(sl, "AI-Powered HR Platform", 1, 2.9, 11, 0.8,
         font_size=28, bold=False, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "Technical Architecture & Project Overview", 1, 3.6, 11, 0.6,
         font_size=18, color=RGBColor(0xCC, 0xCC, 0xCC), align=PP_ALIGN.CENTER)

# Pill badges
pill(sl, "Claude AI",        3.5,  5.2, 1.7, 0.38, TEAL,  WHITE)
pill(sl, "Salesforce",       5.4,  5.2, 1.7, 0.38, RED,   WHITE)
pill(sl, "ServiceNow",       7.3,  5.2, 1.7, 0.38, GREEN, WHITE)
pill(sl, "Next.js 16",       9.2,  5.2, 1.7, 0.38, CHARCOAL, GOLD)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Overview
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Overview", "Two distinct experiences — one Next.js project")

# Left box — Marketing
add_rect(sl, 0.5, 1.6, 5.8, 5.2, WHITE)
add_rect(sl, 0.5, 1.6, 5.8, 0.45, RED)
add_text(sl, "Marketing Site", 0.6, 1.65, 5.6, 0.4,
         font_size=16, bold=True, color=WHITE)
add_text(sl, "(marketing) route group", 0.6, 2.15, 5.5, 0.35,
         font_size=12, color=TEAL, bold=True)
pages = [
    "/  — Home (hero, carousel, benefits)",
    "/products  — All 7 AI agents",
    "/products/[slug]  — Product detail",
    "/use-cases  — Workflow table + industry packs",
    "/templates  — Filterable template library",
    "/company  — About, mission, values",
    "/demo  — Demo request form",
    "/terms  — Terms & conditions",
]
for i, p in enumerate(pages):
    add_text(sl, f"• {p}", 0.65, 2.55 + i * 0.38, 5.5, 0.38,
             font_size=12, color=CHARCOAL)

# Right box — Dashboard
add_rect(sl, 7.0, 1.6, 5.8, 5.2, WHITE)
add_rect(sl, 7.0, 1.6, 5.8, 0.45, TEAL)
add_text(sl, "HR Demo App", 7.1, 1.65, 5.6, 0.4,
         font_size=16, bold=True, color=WHITE)
add_text(sl, "(dashboard) route group", 7.1, 2.15, 5.5, 0.35,
         font_size=12, color=RED, bold=True)
modules = [
    ("/schedule",   "Interview Scheduling",  "Static data"),
    ("/sourcing",   "Candidate Sourcing",    "Static data"),
    ("/onboarding", "Onboarding Assistant",  "ServiceNow + Claude AI"),
    ("/benefits",   "Benefits Assistant",    "ServiceNow + Claude AI"),
    ("/payroll",    "Payroll Assistant",     "Salesforce Agentforce"),
]
for i, (path, name, backend) in enumerate(modules):
    y = 2.55 + i * 0.75
    add_text(sl, path, 7.1, y, 2.2, 0.35, font_size=12, bold=True, color=TEAL)
    add_text(sl, name, 9.4, y, 3.2, 0.35, font_size=12, color=CHARCOAL)
    add_text(sl, backend, 7.1, y + 0.35, 5.5, 0.35, font_size=10,
             color=RGBColor(0x6B, 0x65, 0x60), italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Tech Stack
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Tech Stack", accent=TEAL)

cols   = [3.5, 9.0]
total  = sum(cols)
left   = (13.33 - total) / 2
rows = [
    ("Layer",                   "Technology",                        True,  CHARCOAL, LIGHT_GREY),
    ("Frontend",                "Next.js 16 (App Router), React 19, TypeScript", False, CHARCOAL, WHITE),
    ("Styling",                 "Tailwind CSS v4, shadcn/ui, Framer Motion",     False, CHARCOAL, WHITE),
    ("AI — Onboarding/Benefits","Anthropic Claude claude-sonnet-4-6 via @anthropic-ai/sdk", False, CHARCOAL, WHITE),
    ("AI — Payroll",            "Salesforce Agentforce Agent API",               False, CHARCOAL, WHITE),
    ("HRIS Backend",            "ServiceNow Table REST API (Basic Auth)",        False, CHARCOAL, WHITE),
    ("CRM Backend",             "Salesforce OAuth 2.0 Client Credentials",       False, CHARCOAL, WHITE),
    ("Auth",                    "OAuth 2.0 (Salesforce), HTTP Basic Auth (ServiceNow)", False, CHARCOAL, WHITE),
]
for i, (c1, c2, bold, tc, rc) in enumerate(rows):
    table_row(sl, [c1, c2], left, 1.55 + i * 0.5, cols,
              row_height=0.5, bg_color=rc, text_color=tc,
              font_size=13, bold=bold)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — HR Demo Modules
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "HR Demo Modules", accent=RED)

cols  = [2.5, 1.8, 3.2, 3.8]
left  = 0.5
hdrs  = ["Module", "Status", "AI Engine", "Backend"]
table_row(sl, hdrs, left, 1.55, cols, row_height=0.45,
          bg_color=CHARCOAL, text_color=GOLD, font_size=13, bold=True)

data = [
    ("Onboarding", "Live",        "Claude AI (claude-sonnet-4-6)", "ServiceNow Table API"),
    ("Benefits",   "Live",        "Claude AI (claude-sonnet-4-6)", "ServiceNow Table API"),
    ("Payroll",    "Live",        "Salesforce Agentforce",         "Salesforce Agent API"),
    ("Schedule",   "Live",        "—",                             "Static data"),
    ("Sourcing",   "Live",        "—",                             "Static data"),
    ("Reference",  "Coming Soon", "—",                             "—"),
    ("Engee",      "Coming Soon", "—",                             "—"),
]
for i, row in enumerate(data):
    bg_c = WHITE if i % 2 == 0 else CREAM
    table_row(sl, list(row), left, 2.0 + i * 0.5, cols,
              row_height=0.5, bg_color=bg_c, font_size=13)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Marketing Pages
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Marketing Pages", "Public-facing website built with Framer Motion + Tailwind CSS v4", accent=RED)

pages = [
    ("/",               "Home",          "Hero section, agent carousel, benefits overview, how it works, security"),
    ("/products",       "Products",      "All 7 AI agents grid with video previews, features bento, use cases"),
    ("/products/[slug]","Product Detail","Per-agent deep-dive: intro, stats, how it works, benefits, integrations, workflows"),
    ("/use-cases",      "Use Cases",     "Animated workflow table (HR Ops / Industry Solutions), industry packs, activation steps"),
    ("/templates",      "Templates",     "Filterable template library (9 templates, 7 categories) with setup time & features"),
    ("/company",        "Company",       "About, mission, vision, core values, join team CTA"),
    ("/demo",           "Demo",          "Demo request form"),
    ("/terms",          "Terms",         "Full terms & conditions (9 sections, Canadian jurisdiction)"),
]
cols = [2.2, 1.8, 8.8]
left = 0.3
table_row(sl, ["Path", "Page", "Description"], left, 1.5, cols,
          row_height=0.42, bg_color=CHARCOAL, text_color=GOLD, font_size=12, bold=True)
for i, (path, page, desc) in enumerate(pages):
    bg_c = WHITE if i % 2 == 0 else CREAM
    table_row(sl, [path, page, desc], left, 1.92 + i * 0.48, cols,
              row_height=0.48, bg_color=bg_c, font_size=11)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Onboarding Module Architecture
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Onboarding Module", "ServiceNow + Claude AI — Agentic Tool-Use Loop", accent=TEAL)

# Left: features
add_rect(sl, 0.4, 1.55, 5.6, 5.3, WHITE)
add_text(sl, "Features", 0.5, 1.6, 5.3, 0.4,
         font_size=15, bold=True, color=TEAL)
features = [
    "Live dashboard — all employee onboarding records from ServiceNow",
    "3 sample employees pre-seeded (pending / in-progress / completed)",
    "Task checklist: equipment, access, documents, training, workspace",
    "Status tracking: pending · in_progress · completed · on_hold",
    "AI chat powered by Claude claude-sonnet-4-6",
]
for i, f in enumerate(features):
    add_text(sl, f"• {f}", 0.5, 2.1 + i * 0.6, 5.4, 0.6,
             font_size=12, color=CHARCOAL)

# Right: tools
add_rect(sl, 6.5, 1.55, 6.4, 5.3, CHARCOAL)
add_text(sl, "Claude AI Tools", 6.6, 1.6, 6.0, 0.4,
         font_size=15, bold=True, color=GOLD)
tools = [
    ("get_employee_onboarding",  "Look up employee record by name or ID"),
    ("update_onboarding_task",   "Mark tasks complete / incomplete"),
    ("get_benefit_types",        "Fetch all available benefit plans"),
    ("get_employee_benefits",    "Get employee's benefit enrollment"),
    ("create_it_incident",       "Raise IT support ticket in ServiceNow"),
]
for i, (name, desc) in enumerate(tools):
    add_text(sl, name, 6.6, 2.15 + i * 0.82, 6.0, 0.35,
             font_size=12, bold=True, color=GOLD)
    add_text(sl, desc, 6.6, 2.5 + i * 0.82, 6.0, 0.35,
             font_size=11, color=LIGHT_GREY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Claude AI Agentic Loop
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Claude AI — Agentic Tool-Use Loop", accent=GOLD)

steps = [
    (GOLD,  CHARCOAL, "1", "User Message",           "Employee asks a question or makes a request"),
    (TEAL,  WHITE,    "2", "Claude Reasons",          "claude-sonnet-4-6 decides which tool(s) to call"),
    (RED,   WHITE,    "3", "Tool Execution",          "Next.js API route calls ServiceNow Table REST API"),
    (TEAL,  WHITE,    "4", "Result Returned",         "ServiceNow data returned to Claude as tool_result"),
    (GREEN, CHARCOAL, "5", "Multi-step Reasoning",    "Claude may call additional tools (chained workflow)"),
    (CHARCOAL, GOLD,  "6", "Final Answer",            "Claude responds when stop_reason === 'end_turn'"),
]

x_start = 0.4
w_num   = 0.55
w_title = 2.8
w_desc  = 8.6
for i, (bg_c, tc, num, title, desc) in enumerate(steps):
    y = 1.55 + i * 0.82
    add_rect(sl, x_start,            y, w_num,   0.68, bg_c)
    add_rect(sl, x_start + w_num,    y, w_title, 0.68, WHITE,      LIGHT_GREY, 0.5)
    add_rect(sl, x_start + w_num + w_title, y, w_desc, 0.68, CREAM, LIGHT_GREY, 0.5)
    add_text(sl, num,   x_start,               y + 0.12, w_num,   0.44,
             font_size=20, bold=True, color=tc, align=PP_ALIGN.CENTER)
    add_text(sl, title, x_start + w_num + 0.1, y + 0.12, w_title - 0.2, 0.44,
             font_size=13, bold=True, color=CHARCOAL)
    add_text(sl, desc,  x_start + w_num + w_title + 0.1, y + 0.12, w_desc - 0.2, 0.44,
             font_size=12, color=CHARCOAL)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Payroll Module (Salesforce Agentforce)
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Payroll Module", "Salesforce Agentforce Agent API", accent=GREEN)

# Left
add_rect(sl, 0.4, 1.55, 5.8, 5.3, WHITE)
add_text(sl, "How It Works", 0.55, 1.6, 5.4, 0.4,
         font_size=15, bold=True, color=GREEN)
steps2 = [
    ("1. Start Session",   "POST /agents/{agentId}/sessions",    "OAuth 2.0 Client Credentials token"),
    ("2. Send Message",    "POST /sessions/{sessionId}/messages", "User message forwarded to Agentforce"),
    ("3. Receive Reply",   "Response body parsed",                "Agent reply extracted and returned"),
    ("4. End Session",     "DELETE /sessions/{sessionId}",        "Session cleaned up after conversation"),
]
for i, (title, api, note) in enumerate(steps2):
    y = 2.15 + i * 1.1
    add_text(sl, title, 0.55, y,        5.4, 0.38, font_size=13, bold=True, color=CHARCOAL)
    add_text(sl, api,   0.55, y + 0.35, 5.4, 0.35, font_size=11, color=TEAL,   bold=True)
    add_text(sl, note,  0.55, y + 0.65, 5.4, 0.35, font_size=11, color=RGBColor(0x6B, 0x65, 0x60))

# Right
add_rect(sl, 6.8, 1.55, 6.1, 5.3, CHARCOAL)
add_text(sl, "Authentication", 6.95, 1.6, 5.7, 0.4,
         font_size=15, bold=True, color=GOLD)
auth_lines = [
    "OAuth 2.0 Client Credentials Flow",
    "• No user login required",
    "• Server-to-server token acquisition",
    "• Token cached and refreshed automatically",
    "",
    "Required Environment Variables",
    "• SALESFORCE_CLIENT_ID",
    "• SALESFORCE_CLIENT_SECRET",
    "• SALESFORCE_MY_DOMAIN",
    "• SALESFORCE_AGENT_ID",
    "• SALESFORCE_LOGIN_URL",
]
for i, line in enumerate(auth_lines):
    bold = line.endswith("Flow") or line.endswith("Variables")
    col  = GOLD if bold else (LIGHT_GREY if line.startswith("•") else WHITE)
    add_text(sl, line, 6.95, 2.15 + i * 0.42, 5.7, 0.42,
             font_size=12, bold=bold, color=col)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — ServiceNow Integration
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "ServiceNow Integration", "Custom Scoped App — x_1926120_employee", accent=TEAL)

# Tables
add_text(sl, "Custom Tables", 0.5, 1.55, 12, 0.4,
         font_size=15, bold=True, color=TEAL)
t_cols = [3.8, 3.8, 5.1]
table_row(sl, ["Table Name", "Purpose", "Key Fields"], 0.5, 1.95, t_cols,
          bg_color=CHARCOAL, text_color=GOLD, font_size=12, bold=True)
sn_tables = [
    ("x_1926120_employee_onboarding",        "Onboarding records & task tracking",  "status, tasks, employee_name, start_date"),
    ("x_1926120_employee_benefit_types",      "Available benefit plan definitions",  "plan_name, type, coverage, cost"),
    ("x_1926120_employee_benefit_enrollment", "Employee benefit enrollments",        "employee_id, plan_id, enrollment_date"),
]
for i, row in enumerate(sn_tables):
    bg_c = WHITE if i % 2 == 0 else CREAM
    table_row(sl, list(row), 0.5, 2.4 + i * 0.5, t_cols,
              bg_color=bg_c, font_size=11)

# Config requirements
add_text(sl, "Required Configuration", 0.5, 4.05, 12, 0.4,
         font_size=15, bold=True, color=TEAL)
reqs = [
    ("Table Settings (System Definition > Tables)",
     "• 'Allow access to this table via web services' — must be checked\n"
     "• 'Accessible from' — set to 'All application scopes'"),
    ("Admin User Roles",
     "• x_1926120_employee.hr_manager\n• x_1926120_employee.hr_representative"),
    ("Authentication",
     "• HTTP Basic Auth — credentials Base64-encoded in Authorization header"),
]
for i, (title, body) in enumerate(reqs):
    x = 0.5 + i * 4.3
    add_rect(sl, x, 4.5, 4.1, 2.6, WHITE)
    add_text(sl, title, x + 0.1, 4.55, 3.9, 0.4, font_size=12, bold=True, color=CHARCOAL)
    add_text(sl, body,  x + 0.1, 5.0,  3.9, 2.0, font_size=11, color=CHARCOAL)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — API Routes
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "API Routes", accent=RED)

cols = [3.8, 1.4, 7.6]
left = 0.3
table_row(sl, ["Route", "Method", "Description"], left, 1.55, cols,
          bg_color=CHARCOAL, text_color=GOLD, font_size=12, bold=True)
routes = [
    ("/api/onboarding/records",           "GET",    "Fetch all onboarding records from ServiceNow"),
    ("/api/onboarding/chat",              "POST",   "Claude AI agentic chat with ServiceNow tool use"),
    ("/api/benefits/records",             "GET",    "Fetch benefit types and enrollments from ServiceNow"),
    ("/api/benefits/chat",                "POST",   "Claude AI agentic chat for benefits queries"),
    ("/api/agent",                        "POST",   "Salesforce Agentforce proxy — start / message / end session"),
    ("/api/auth/salesforce",              "GET",    "Initiate Salesforce OAuth 2.0 flow"),
    ("/api/auth/callback/salesforce",     "GET",    "Salesforce OAuth callback — exchange code for token"),
    ("/api/auth/status",                  "GET",    "Check current Salesforce auth status"),
]
method_colors = {"GET": GREEN, "POST": RED, "DELETE": RGBColor(0xCC, 0x44, 0x44)}
for i, (route, method, desc) in enumerate(routes):
    bg_c = WHITE if i % 2 == 0 else CREAM
    x = left
    # Route cell
    add_rect(sl, x, 2.0 + i * 0.48, cols[0], 0.48, bg_c, LIGHT_GREY, 0.5)
    add_text(sl, route, x + 0.08, 2.06 + i * 0.48, cols[0] - 0.16, 0.38,
             font_size=11, bold=True, color=TEAL)
    x += cols[0]
    # Method cell
    mc = method_colors.get(method, CHARCOAL)
    add_rect(sl, x, 2.0 + i * 0.48, cols[1], 0.48, bg_c, LIGHT_GREY, 0.5)
    add_text(sl, method, x + 0.08, 2.06 + i * 0.48, cols[1] - 0.16, 0.38,
             font_size=11, bold=True, color=mc)
    x += cols[1]
    # Desc cell
    add_rect(sl, x, 2.0 + i * 0.48, cols[2], 0.48, bg_c, LIGHT_GREY, 0.5)
    add_text(sl, desc, x + 0.08, 2.06 + i * 0.48, cols[2] - 0.16, 0.38,
             font_size=11, color=CHARCOAL)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Project Structure
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Project Structure", accent=CHARCOAL)

# Left column
add_rect(sl, 0.4, 1.55, 6.1, 5.6, CHARCOAL)
structure_left = """\
src/app/
  (marketing)/
    page.tsx            ← Home /
    products/
      page.tsx          ← /products
      [slug]/page.tsx   ← /products/:slug
    use-cases/page.tsx
    templates/page.tsx
    company/page.tsx
    demo/page.tsx
    terms/page.tsx
  (dashboard)/
    schedule/
    sourcing/
    onboarding/
    benefits/
    payroll/
  api/
    agent/
    onboarding/chat|records
    benefits/chat|records"""
add_text(sl, structure_left, 0.55, 1.65, 5.8, 5.3,
         font_size=10, color=GOLD, bold=False)

# Right column
add_rect(sl, 6.8, 1.55, 6.1, 5.6, WHITE)
structure_right = """\
src/components/
  marketing/
    Layout/     ← Header + Footer
    Home/       ← 7 sections
    Products/   ← 7 sections
    UseCases/   ← 6 sections
    Templates/  ← 5 sections
    Company/    ← 5 sections
    Terms/
  onboarding|benefits|payroll/
  layout/       ← Sidebar + Header
  ui/           ← shadcn + custom
  ProductIntroduction.tsx
  ProductHowItWorks.tsx
  ProductBenefits.tsx
  hero-with-mockup.tsx
  bento-grid_2.tsx

src/data/
  products.ts   ← 7 agents full data

src/lib/
  servicenow.ts
  salesforce.ts
  modules.ts"""
add_text(sl, structure_right, 6.95, 1.65, 5.8, 5.3,
         font_size=10, color=CHARCOAL)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Getting Started
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Getting Started", accent=GREEN)

steps3 = [
    (GREEN,    "1",  "Install Dependencies",
               "cd haigent-project && npm install"),
    (TEAL,     "2",  "Copy Public Assets",
               "Copy /animations/, /svg_I/, /all_robo.png, /hero_image*.png,\n/Logo_simple_black.png from haigent.ai/public/ → haigent-project/public/"),
    (GOLD,     "3",  "Configure .env.local",
               "SALESFORCE_CLIENT_ID / CLIENT_SECRET / MY_DOMAIN / AGENT_ID\nSERVICENOW_INSTANCE_URL / USERNAME / PASSWORD\nANTHROPIC_API_KEY"),
    (RED,      "4",  "Delete Route Conflict",
               "Delete src/app/page.tsx — the marketing (marketing)/page.tsx\nserves / instead; having both causes a Next.js build error"),
    (CHARCOAL, "5",  "Run Dev Server",
               "npm run dev   →   http://localhost:3000\nMarketing home at /   |   HR demo dashboard at /schedule"),
]

for i, (color, num, title, body) in enumerate(steps3):
    y = 1.55 + i * 1.05
    add_rect(sl, 0.4, y, 0.55, 0.85, color)
    add_text(sl, num, 0.4, y + 0.17, 0.55, 0.55,
             font_size=22, bold=True, color=WHITE if color != GOLD else CHARCOAL,
             align=PP_ALIGN.CENTER)
    add_rect(sl, 0.95, y, 11.95, 0.85, WHITE, LIGHT_GREY, 0.5)
    add_text(sl, title, 1.1, y + 0.05, 11.6, 0.38,
             font_size=13, bold=True, color=CHARCOAL)
    add_text(sl, body, 1.1, y + 0.45, 11.6, 0.5,
             font_size=11, color=RGBColor(0x6B, 0x65, 0x60))


# ── Save ──────────────────────────────────────────────────────────────────────
output = "Haigent_Project.pptx"
prs.save(output)
print(f"✓  Saved → {output}  ({prs.slides.__len__()} slides)")
